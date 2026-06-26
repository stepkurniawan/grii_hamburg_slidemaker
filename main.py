"""
This program is to make a powerpoint presentation for Sunday service.
It is tweaked for Hamburg Church, which is bilingual (Indonesian and English).
The inputs that are changed every week:
1. 4 song's numbers 
2. Opening Bible verse
3. Pastor's title in Indonesian and English, and name
4. Second (blue) offering purpose in Indonesian and English

The songs are stored in google drive as a jpg. 
So we just have to use that picture as a slide. 

The slide will be generated based on this structure:
1. Our service is about to start page
2. church cover - page (Static)
3. First song
4. church cover - page
5. Second song
6. Opening Bible vers (in English), ex: Subtitel: Scripture Reading, Titel: Jeremiah 20:1-2 
    - English Translation with the verse number
    - Indonesian Translation with the verse number
7. church cover - page
8. Third song
9. church cover - page
10. Lord's Prayer (in English and Indonesian) -  3 slides - Static, we have it
11. church cover - page
12. Predigt - page
    - Title: <Pastor's title in English>. <Pastor's name>
    - Subtitle: <Pastor's title in Indonesian>. <Pastor's name>
13. church cover - page
14. Apostles' Creed (in English and Indonesian) - 6 slides - Static, we have it
15. church cover - page
16. Offerings - page
    - Small Title: Kollekte
    - Title: MRII Hamburg (rot) & <Second offering purpose in English> (blau)
    - SubTitle: MRII Hamburg (merah) & <Second offering purpose in Indonesian> (biru)
17. Fourth song
18. church cover - page
19. Puji Allah Bapa Putra - page - 3 slides - Static, we have it
20. church cover - page
21. A...Men - page - 1 slide - Static, we have it
22. Annoucement - page - 1 slide - Static
23. Persekutuan Doa - page - 1 slide - Static
24. Seminar - page - 1 slide - Static
25. Ibadah Minggu Depan - page - 1 slide - Static
26. Sekolah Minggu - page - 1 slide - Static
27. Happy Birthday - page - 1 slide - Static
28. Coffee Time - page - 1 slide - Static
    
"""

####### Versioning
# 3.0.0 implement in memory song download
# 3.1.0 implement english if german cant be found
# "3.1.2 bug fix english path
# 3.2.0 implement multi ayat alkitab
# 3.2.1 bug fix Pengkothbah typo, add better bible verse cover
# 3.2.2 make bible verse more robust
# 3.2.3 change the bible input from indo to english
# 3.4.0 refactor adding pydantic 

# Importing libraries

import datetime
from io import BytesIO
import os
import sys
import streamlit as st

from pydantic import ValidationError
from pptx import Presentation
from pptx.util import Inches

from annoucement import insert_annoucement_slides
from models import ServiceOrder
from Pujian import SONGS_FOLDER
from footer import footer
from pptx_creator import (
    add_amen_page,
    add_appostle_creed_page,
    add_beginning_slide,
    add_bekantmachung_page,
    add_bible_reading_page,
    add_church_cover_page,
    add_doa_bapa_kami_page,
    add_doxology_page,
    add_intersession_page,
    add_preacher_page,
    add_secondary_offering_purpose_page,
    decide_offering_purpose_layout_name,
    insert_song_slides_drive_folder,
)


########################################### Version ##########################################
# date today in yymmdd
date_today = datetime.datetime.now().strftime("%y%m%d")
VERSION = "3.3."+str(date_today)

########################################### INPUTS ##########################################
global SONG_NUMBERS, PASTOR_TITLE_ID, PASTOR_NAME, OPEN_BIBLE_FULL_VERSE,OPEN_BIBLE_FULL_VERSES , PASTOR_TITLE_DE_OR_EN, SELECTED_SECOND_OFFERING_PURPOSE_ID

# 1. 4 song's numbers
SONG_NUMBERS = ["1", "test", "test", "test"]
# 2. Opening Bible verse
OPEN_BIBLE_FULL_VERSES = ["Genesis 1:2-3", "1 Kings 1:1-2"]
OPEN_BIBLE_FULL_VERSE = "Genesis 1:2-3"
# 3. Pastor's title in Indonesian and English, and name
PASTOR_TITLE_ID = "Pdt."
PASTOR_TITLE_DE_OR_EN = "Rev."
PASTOR_NAME = "Billy Kristanto"
# 4. Second (blue) offering purpose in Indonesian and English
SECOND_OFFERING_PURPOSE_ID = ["NONE", "P_PENGINJILAN", "P_SEKOLAH", "P_MANDAT", "P_PEMBANGUNAN", "P_DIAKONIA" ]
SELECTED_SECOND_OFFERING_PURPOSE_ID = "NONE"

########################################### Global variables##########################################
MY_SLIDE_WIDTH = Inches(16)
MY_SIDE_HEIGHT = Inches(9)
CURRENT_DIR = getattr(sys, '_MEIPASS', os.path.dirname(os.path.abspath(__file__)))

HOME_DIR = os.path.expanduser("~")
DOWNLOAD_FOLDER = os.path.join(CURRENT_DIR, "Downloads")
# print("HOME_DIR: " + HOME_DIR)
GRII_FOLDER = os.path.join(DOWNLOAD_FOLDER, "GRII")
OUTPUT_DIR = os.path.join(GRII_FOLDER,"GRII_Slides")
NEW_OUTPUT_DIR = os.path.join(CURRENT_DIR,"Output")
output_file = ""
binary_output_file = BytesIO()


########################################### Functions ##########################################

def st_print(text):
    st.write(text)
    print(text)

def st_error_print(text):
    st.error(text)
    print(text)

def get_resource_path(relative_path):
    if hasattr(sys, "_MEIPASS"):
        # Running from PyInstaller executable, use sys._MEIPASS
        base_path = sys._MEIPASS
    else:
        # Running from source code
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)

TEMPLATE_FILE = get_resource_path('master_slide_template_en.pptx')

def processing_answers(data_array):
    # Legacy adapter for callers that still pass the old positional list.
    return ServiceOrder.model_validate(
        {
            "song_numbers": data_array[0],
            "pastor_name": data_array[1],
            "bible_verses": data_array[2],
            "pastor_title": data_array[3],
            "holy_communion_song_number": data_array[4],
        }
    )


def default_service_order():
    return ServiceOrder.model_validate(
        {
            "song_numbers": "161, 320, 93, 169",
            "pastor_name": "Pdt. Billy Kristanto",
            "bible_verses": "Genesis 1:2-3, 1 Kings 1:1-2",
            "pastor_title": "Rev.",
            "holy_communion_song_number": None,
        }
    )


def show_validation_errors(error: ValidationError):
    for validation_error in error.errors():
        location = " > ".join(str(part) for part in validation_error["loc"])
        message = validation_error["msg"]
        st.sidebar.error(f"{location}: {message}")

def sunday_date(formatted):
    # save file as with next sunday's date yyyymmdd.pptx
    # Get today's date
    today = datetime.date.today()
    # Calculate the number of days until the next Sunday (0 = Sunday, 1 = Monday, ..., 6 = Saturday)
    days_until_sunday = (6 - today.weekday()) % 7
    # Calculate the date of the next Sunday
    next_sunday = today + datetime.timedelta(days=days_until_sunday)
    if formatted == "filename":
        next_sunday = str(next_sunday).replace("-", "")
    elif formatted == "slide":
        # return 12 July 2020
        next_sunday = str(next_sunday.strftime("%d %B %Y"))
    else:
        next_sunday = next_sunday
        
    return next_sunday

def create_website():
    # website title
    page_title = "☁️ MRII Europe SlideMaker"

    st.set_page_config( page_icon=":church:", page_title=page_title)

    # website description
    # heading
    st.title(page_title)
    # subheading
    st.subheader(" 📜 Welcome to MRII Europe SlideMaker V" + VERSION + " A1")
    st.write("This website is used to create a powerpoint presentation for Sunday service.")
    st.write("It is tweaked for Hamburg Church, which is bilingual (English and Indo).")

    # subheading
    st.subheader("How to use this website")
    st.write("1. Enter the song numbers separated by comma on the left sidebar")
    st.write("2. Enter the (English) Bible verse. ex: Genesis 1:2-3, 1 Kings 1:1-2")
    st.write("3. (Optional) Enter the pastor name ")
    st.write("4. (Optional) Enter the pastor title in English ")
    st.write("5. Click the submit button")
    
    # subheading
    st.subheader("Where are the songs?")
    st.write("The songs are downloaded from the database to the in-memory, since we don't have permission to store it in the server.")
    st.write("It is stored in the folder, in the server: " + SONGS_FOLDER)

    st.subheader("Where is the Slide?")
    st.markdown("After the main process is finished, the **Download Button** will appear on the **left sidebar**")
    st.markdown("You can save the slide by clicking on the **Download Button**")
    st.markdown("The file name is the date of the **next Sunday** service")

    # ask for user input as parameters on the side bar
    # ask for song numbers
    st.sidebar.subheader("Song numbers")
    st.sidebar.write("Please enter the song numbers separated by comma")
    st.sidebar.write("Example: 161, 320, 93, 169")
    song_numbers = st.sidebar.text_input("Song numbers")

    holy_communion = st.sidebar.checkbox("Holy Communion", key="holy_communion", help="Toggle this if the service is Holy Communion. This will add the Holy Communion song slide before the sermon.")
    holy_communion_song_number: int | None = None
    if holy_communion:
        st.sidebar.write("Holy Communion song will be added before the sermon slide.")
        holy_communion_song_number = st.sidebar.text_input("Holy Communion song number")
        if holy_communion and not holy_communion_song_number:
            st.sidebar.error("Please enter the Holy Communion song number.")

    # ask for Bible verse
    st.sidebar.subheader("Bible verse")
    st.sidebar.write("Please enter the Bible verse in English")
    st.sidebar.write("Example: Genesis 1:2-3, 1 Kings 1:1-2")
    bible_verses = st.sidebar.text_input("Bible verse(s)")
    
    # ask for pastor name
    st.sidebar.subheader("Pastor name")
    st.sidebar.write("Please enter the pastor name")
    st.sidebar.write("Default Example: Pdt. Billy Kristanto")
    pastor_name = st.sidebar.text_input("Pastor name")

    # if pastor_name is empty, then use default value
    if pastor_name == "":
        pastor_name = "Pdt. Billy Kristanto"

    # ask for pastor title in English
    st.sidebar.subheader("Pastor title")
    st.sidebar.write("Please enter the pastor title")
    st.sidebar.write("Default Example: Rev.")
    # default is Rev. 
    pastor_title = st.sidebar.text_input("Pastor title")

    # if pastor_title_de is empty, then use default value
    if pastor_title == "":
        pastor_title = "Rev."


    # create a submit button
    submit_button = st.sidebar.button("Submit")
    # if submit button is clicked
    if submit_button:
        try:
            service_order = ServiceOrder.model_validate(
                {
                    "song_numbers": song_numbers,
                    "pastor_name": pastor_name,
                    "bible_verses": bible_verses,
                    "pastor_title": pastor_title,
                    "holy_communion_song_number": holy_communion_song_number,
                }
            )
        except ValidationError as error:
            st.sidebar.error("Please fix the input before generating the slide.")
            show_validation_errors(error)
            footer()
            return

        st.write("Validated service order")
        st.write(service_order.model_dump(mode="json"))
        with st.spinner('Generating the slide...'):
            main(service_order)
            st.balloons()
            st.sidebar.success("Slide generated successfully in " + NEW_OUTPUT_DIR + ":tada:" + "https://drive.google.com/drive/folders/1AJTLk-AXOI7nEYcWAOMTZ_MxDNzaRSK2")

            st.sidebar.download_button(
                label="Download slide!",
                data=binary_output_file.getvalue(),
                file_name=sunday_date("filename")+ ".pptx",
            )
    
    footer()


def main(service_order: ServiceOrder | None = None):

    """
    the structure of the slide is
    Presentation -> Layout name -> slide layout -> slide -> shapes -> placeholders
    """

    if service_order is None:
        service_order = default_service_order()

    binary_output_file.seek(0)
    binary_output_file.truncate(0)

    ########################################### CHECKING SONGS ##########################################

    # create a test presentation file
    prs = Presentation(TEMPLATE_FILE)

    #### Slide creation starts here ####
    add_beginning_slide(prs)
    add_church_cover_page(prs, sunday_date("slide"))
    # check_placeholders_in_slide_index(prs, 4)

    #### add first song
    st_print("Adding first song")
    try:
        insert_song_slides_drive_folder(prs, service_order.songs.worship_songs[0].value)
    except Exception:
        print("Error: Cannot add the 1st song")
        st_error_print("Error: Cannot add the 1st song")

    add_church_cover_page(prs, sunday_date("slide"))

    #### add second song
    st_print("Adding second song")
    try: 
        insert_song_slides_drive_folder(prs, service_order.songs.worship_songs[1].value)
    except Exception:
        print("Error: Cannot add the 2nd song")
        st_error_print("Error: Cannot add the 2nd song")

    add_church_cover_page(prs, sunday_date("slide"))

    #### ADD BIBLE VERSE
    st_print("Adding bible reading")

    print("Bible references: ", service_order.bible_references)
    for bible_verse in service_order.bible_references:
        print("bible_verse: ", bible_verse)
        add_bible_reading_page(prs, bible_verse.as_reference_text())

    add_church_cover_page(prs, sunday_date("slide"))

    #### add third song
    st_print("Adding third song")
    try: 
        insert_song_slides_drive_folder(prs, service_order.songs.worship_songs[2].value)
    except Exception:
        print("Error: Cannot add the 3rd song")
        st_error_print("Error: Cannot add the 3rd song")

    add_church_cover_page(prs, sunday_date("slide"))

    
    st_print("Adding Lord's Prayer")
    add_intersession_page(prs)
    add_doa_bapa_kami_page(prs)

    add_church_cover_page(prs, sunday_date("slide"))

    if service_order.songs.holy_communion_song is not None:
        st_print("Adding Holy Communion song")
        try:
            insert_song_slides_drive_folder(prs, service_order.songs.holy_communion_song.value)
            add_church_cover_page(prs, sunday_date("slide"))
        except Exception:
            print("Error: Cannot add the Holy Communion song")
            st_error_print("Error: Cannot add the Holy Communion song")

    st_print("Adding Preacher Sermon Page")
    add_preacher_page(
        prs,
        service_order.pastor.title_id,
        service_order.pastor.title_de_or_en,
        service_order.pastor.name,
    )

    add_church_cover_page(prs, sunday_date("slide"))

    st_print("Adding Apostles' Creed")
    add_appostle_creed_page(prs)

    add_church_cover_page(prs, sunday_date("slide"))

    st_print("Adding Offerings")
    secondary_purpose_id = decide_offering_purpose_layout_name(sunday_date("date"))
    add_secondary_offering_purpose_page(prs, secondary_purpose_id) 

    #### add fourth song
    st_print("Adding fourth song")
    try:
        insert_song_slides_drive_folder(prs, service_order.songs.worship_songs[3].value)
    except Exception:
        print("Error: Cannot add the 4th song")
        st_error_print("Error: Cannot add the 4th song")

    add_church_cover_page(prs, sunday_date("slide"))

    st_print("Adding Puji Allah Bapa Putra until closings")
    add_doxology_page(prs)

    add_church_cover_page(prs, sunday_date("slide"))

    add_amen_page(prs)

    add_church_cover_page(prs, sunday_date("slide"))

    add_bekantmachung_page(prs)
    insert_annoucement_slides(prs)
    
    # save in output folder
    prs.save(binary_output_file) # so it is downloadable using button
    st_print("Saving the slide in " + NEW_OUTPUT_DIR)
    prs.save(os.path.join(NEW_OUTPUT_DIR, sunday_date("filename") + ".pptx"))
    st_print("saved in " + NEW_OUTPUT_DIR)


# Run the Streamlit UI when this file is used as the app entry point.
if __name__ == "__main__":
    create_website()


