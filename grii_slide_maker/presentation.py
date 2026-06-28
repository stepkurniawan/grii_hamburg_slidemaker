"""PowerPoint assembly flow for the Sunday service slide deck."""

from collections.abc import Callable
from dataclasses import dataclass
from io import BytesIO
import os

from pptx import Presentation

from grii_slide_maker.models import OrderOfMass
from grii_slide_maker.slides.announcements import insert_annoucement_slides
from grii_slide_maker.slides.creator import (
    add_amen_page,
    add_appostle_creed_page,
    add_beginning_slide,
    add_bekantmachung_page,
    add_bible_reading_page,
    add_church_cover_page,
    add_doa_bapa_kami_page,
    add_doxology_page,
    add_intersession_page,
    add_preacher_page,
    add_secondary_offering_purpose_page,
    decide_offering_purpose_layout_name,
    insert_song_slides_drive_folder,
)


@dataclass(frozen=True)
class SlideDeckActions:
    add_amen_page: Callable = add_amen_page
    add_appostle_creed_page: Callable = add_appostle_creed_page
    add_beginning_slide: Callable = add_beginning_slide
    add_bekantmachung_page: Callable = add_bekantmachung_page
    add_bible_reading_page: Callable = add_bible_reading_page
    add_church_cover_page: Callable = add_church_cover_page
    add_doa_bapa_kami_page: Callable = add_doa_bapa_kami_page
    add_doxology_page: Callable = add_doxology_page
    add_intersession_page: Callable = add_intersession_page
    add_preacher_page: Callable = add_preacher_page
    add_secondary_offering_purpose_page: Callable = add_secondary_offering_purpose_page
    decide_offering_purpose_layout_name: Callable = decide_offering_purpose_layout_name
    insert_annoucement_slides: Callable = insert_annoucement_slides
    insert_song_slides_drive_folder: Callable = insert_song_slides_drive_folder


def build_service_slides(
    order_of_mass: OrderOfMass,
    *,
    template_file: str,
    output_dir: str,
    binary_output_file: BytesIO,
    sunday_date: Callable[[str], object],
    status_writer: Callable[[str], None],
    error_writer: Callable[[str], None],
    actions: SlideDeckActions | None = None,
) -> None:
    """
    Build the Sunday service presentation and save it to memory and disk.

    The structure of the slide is:
    Presentation -> Layout name -> slide layout -> slide -> shapes -> placeholders.
    """
    actions = actions or SlideDeckActions()
    binary_output_file.seek(0)
    binary_output_file.truncate(0)

    prs = Presentation(template_file)

    actions.add_beginning_slide(prs)
    actions.add_church_cover_page(prs, sunday_date("slide"))

    add_song_with_status(
        prs,
        order_of_mass.songs.worship_songs[0].value,
        status="Adding first song",
        error="Error: Cannot add the 1st song",
        status_writer=status_writer,
        error_writer=error_writer,
        actions=actions,
    )

    actions.add_church_cover_page(prs, sunday_date("slide"))

    add_song_with_status(
        prs,
        order_of_mass.songs.worship_songs[1].value,
        status="Adding second song",
        error="Error: Cannot add the 2nd song",
        status_writer=status_writer,
        error_writer=error_writer,
        actions=actions,
    )

    actions.add_church_cover_page(prs, sunday_date("slide"))

    status_writer("Adding bible reading")
    print("Bible references: ", order_of_mass.bible_references)
    for bible_verse in order_of_mass.bible_references:
        print("bible_verse: ", bible_verse)
        add_bible_reading_with_status(
            prs,
            bible_verse.as_reference_text(),
            error_writer=error_writer,
            actions=actions,
        )

    actions.add_church_cover_page(prs, sunday_date("slide"))

    add_song_with_status(
        prs,
        order_of_mass.songs.worship_songs[2].value,
        status="Adding third song",
        error="Error: Cannot add the 3rd song",
        status_writer=status_writer,
        error_writer=error_writer,
        actions=actions,
    )

    actions.add_church_cover_page(prs, sunday_date("slide"))

    status_writer("Adding Lord's Prayer")
    actions.add_intersession_page(prs)
    actions.add_doa_bapa_kami_page(prs)

    actions.add_church_cover_page(prs, sunday_date("slide"))

    if order_of_mass.songs.holy_communion_song is not None:
        status_writer("Adding Holy Communion song")
        try:
            actions.insert_song_slides_drive_folder(
                prs,
                order_of_mass.songs.holy_communion_song.value,
            )
            actions.add_church_cover_page(prs, sunday_date("slide"))
        except Exception:
            print("Error: Cannot add the Holy Communion song")
            error_writer("Error: Cannot add the Holy Communion song")

    status_writer("Adding Preacher Sermon Page")
    actions.add_preacher_page(
        prs,
        order_of_mass.pastor.title_id,
        order_of_mass.pastor.title_de_or_en,
        order_of_mass.pastor.name,
    )

    actions.add_church_cover_page(prs, sunday_date("slide"))

    status_writer("Adding Apostles' Creed")
    actions.add_appostle_creed_page(prs)

    actions.add_church_cover_page(prs, sunday_date("slide"))

    status_writer("Adding Offerings")
    secondary_purpose_id = actions.decide_offering_purpose_layout_name(sunday_date("date"))
    actions.add_secondary_offering_purpose_page(prs, secondary_purpose_id)

    add_song_with_status(
        prs,
        order_of_mass.songs.worship_songs[3].value,
        status="Adding fourth song",
        error="Error: Cannot add the 4th song",
        status_writer=status_writer,
        error_writer=error_writer,
        actions=actions,
    )

    actions.add_church_cover_page(prs, sunday_date("slide"))

    status_writer("Adding Puji Allah Bapa Putra until closings")
    actions.add_doxology_page(prs)

    actions.add_church_cover_page(prs, sunday_date("slide"))

    actions.add_amen_page(prs)

    actions.add_church_cover_page(prs, sunday_date("slide"))

    actions.add_bekantmachung_page(prs)
    actions.insert_annoucement_slides(prs)

    os.makedirs(output_dir, exist_ok=True)
    prs.save(binary_output_file)
    status_writer("Saving the slide in " + output_dir)
    output_path = os.path.join(output_dir, sunday_date("filename") + ".pptx")
    with open(output_path, "wb") as saved_pptx:
        saved_pptx.write(binary_output_file.getvalue())
    status_writer("saved in " + output_dir)


def add_song_with_status(
    prs: object,
    song_number: str,
    *,
    status: str,
    error: str,
    status_writer: Callable[[str], None],
    error_writer: Callable[[str], None],
    actions: SlideDeckActions,
) -> None:
    status_writer(status)
    try:
        actions.insert_song_slides_drive_folder(prs, song_number)
    except Exception:
        print(error)
        error_writer(error)


def add_bible_reading_with_status(
    prs: object,
    bible_reference: str,
    *,
    error_writer: Callable[[str], None],
    actions: SlideDeckActions,
) -> None:
    try:
        actions.add_bible_reading_page(prs, bible_reference)
    except Exception as error:
        message = f"Error: Cannot add Bible reading {bible_reference}: {error}"
        print(message)
        error_writer(message)
