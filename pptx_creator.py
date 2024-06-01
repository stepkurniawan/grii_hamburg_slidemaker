import os
import streamlit as st
from PIL import Image

# from alkitab_scraper import get_ayat_alkitab_dict
from pptx.util import Inches
from bible_translation import indonesian_to_german_bible
from bible_translation import lai_abbre_to_full
from Bible_API import get_verses_dict
from Pujian import *


def st_print(text):
    st.write(text)
    print(text)

def sort_by_number(file_name):
    # Custom sorting function to extract numbers from the file name and sort numerically
    try:
        key = int(''.join(filter(str.isdigit, file_name)))
    except ValueError:
        key = 99
        st_print("Error: file name does not contain any number: ", file_name)
    return key

# create slides from this folder. one slide for each file. The folder contains jpg files, and it should be scaled to fit the slide. 
def insert_slides_from_pict_folder(prs, folder_path):
    # get all files in the folder
    files = sorted(os.listdir(folder_path), key=sort_by_number)

    # loop through all files (picture only) can be png, jpg, jpeg, etc, or it can also be uppercase
    # it has to be sorted by name
    for file in files:
        # check if the file is a picture
        if file.lower().endswith((".jpg", ".png", ".jpeg")):
            # create a new slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # add the picture to the slide
            pic = slide.shapes.add_picture(os.path.join(folder_path, file), Inches(0), Inches(0), height=prs.slide_height, width=prs.slide_width)

def add_slide_layout_from_layout_name(prs,layout_name):
    # Specify the layout name you want to use
    slide = None
    slide_layout = None
    for layout in prs.slide_master.slide_layouts:
        if layout.name == layout_name:
            slide_layout = layout
            break

    # Check if the layout was found
    if slide_layout is not None:
        # Add a slide based on the layout
        slide = prs.slides.add_slide(slide_layout)

    return slide

def add_beginning_slide(prs):
    # Specify the layout name you want to use
    layout_name = "beginning" # renamed in the master template pptx file
    add_slide_layout_from_layout_name(prs, layout_name)



def add_church_cover_page(prs, sunday_date):
    # Specify the layout name you want to use
    layout_name = "COVER_2" # renamed in the master template pptx file
    
    # print("cover layout")
    slide_layout = add_slide_layout_from_layout_name(prs, layout_name)
    # check_placeholders_in_slide(prs,slide_layout)

    sundays_date_placeholder = slide_layout.placeholders[10]
    sundays_date_placeholder.text = sunday_date
    

def get_full_book_name(book_name):
    # if bible_book has less than 4 characters, then its the abbreviated version. the full version is in lai_abbre_to_full dictionary
    output = book_name
    if len(book_name) < 4:
        output = lai_abbre_to_full[book_name]
    # if bible_book has 4 characters, but the first one is a number, then its the abbreviated version. the full version is in lai_abbre_to_full dictionary
    elif book_name[0].isdigit() & (len(book_name) == 4):
        output = lai_abbre_to_full[book_name]
    
    return output

def add_bible_reading_page(prs, bible_verse_text = "Kej 1:2-3"):
    # clean bible_verse_text 
    bible_verse_text = bible_verse_text.replace(".", "")
    # capitalize the first letter if its not number like 1Sam
    if not bible_verse_text[0].isdigit():
        bible_verse_text = bible_verse_text[0].upper() + bible_verse_text[1:]
    elif bible_verse_text[0].isdigit() & (len(bible_verse_text) == 4):
        # capitalize the second letter, ex: 1sam -> 1Sam
        bible_verse_text = bible_verse_text[0] + bible_verse_text[1].upper() + bible_verse_text[2:]

        
    
    
    # Find the layout with the specified name
    LAYOUT_NAME_COVER = "BIBLE_READING" # Bible reading cover
    slide_layout_cover = add_slide_layout_from_layout_name(prs, LAYOUT_NAME_COVER)
    bible_book_ID = ""
    bible_book_DE = ""

    print("bible verse text: ", bible_verse_text)
    bible_book = bible_verse_text.split(" ")[0] # 2Sam
    print("bible_book: ", bible_book)
    bible_book_ID = get_full_book_name(bible_book) # 2 Samuel

    bible_chapter = bible_verse_text.split(" ")[1].split(":")[0]
    bible_verse_start = bible_verse_text.split(" ")[1].split(":")[1].split("-")[0]
    bible_verse_end = bible_verse_text.split(" ")[1].split(":")[1].split("-")[1]

    
    ## ADD BIBLE RADING COVER PAGE
    try: 
        # check_placeholders_in_slide(prs,slide_layout_cover)
        bible_verse = slide_layout_cover.placeholders[10]
        bible_cover_text = bible_book_ID + " " + bible_chapter + ":" + bible_verse_start + "-" + bible_verse_end # Kejadian 1:2-3
        bible_verse.text = bible_cover_text 

        bible_verse_DE = slide_layout_cover.placeholders[11]
        bible_book_DE = indonesian_to_german_bible.get(bible_book_ID)
        bible_cover_text_DE = bible_book_DE + " " + bible_chapter + ":" + bible_verse_start + "-" + bible_verse_end # Genesis 1:2-3
        bible_verse_DE.text = bible_cover_text_DE


    except IndexError:
        print("Invalid placeholder index.")



    # get the bible verse in german
    bible_book_DE = indonesian_to_german_bible.get(bible_book_ID)

    id_bible_verse = get_verses_dict(bible_book_ID, bible_chapter, bible_verse_start, bible_verse_end, "ID") 
    de_bible_verse = get_verses_dict(bible_book_ID, bible_chapter, bible_verse_start, bible_verse_end, "DE")

    # count how many verse 
    verse_count = len(id_bible_verse)

    # loop so we add slide for each verse
    for id_key, de_key in zip(id_bible_verse.keys(), de_bible_verse.keys()):
        bible_verse_layout_name = "BIBLE_VERSE" # renamed in the master template pptx file
        slide_layout_bible_verse = add_slide_layout_from_layout_name(prs, bible_verse_layout_name)
        try:
            # print("bible verse layout")
            # check_placeholders_in_slide(prs,slide_layout_bible_verse)
            de_bible_verse_placeholder = slide_layout_bible_verse.placeholders[10]
            id_bible_verse_placeholder = slide_layout_bible_verse.placeholders[11]

            # Get the values corresponding to the current keys
            de_value = de_bible_verse[de_key]
            id_value = id_bible_verse[id_key]

            de_bible_verse_placeholder.text = de_key + " : " + de_value
            id_bible_verse_placeholder.text = id_key + " : " + id_value
            
        except IndexError:
            print("Invalid placeholder index.")
    



def add_doa_bapa_kami_page(prs):
    add_slide_layout_from_layout_name(prs, "BAPA_KAMI_1")
    add_slide_layout_from_layout_name(prs, "BAPA_KAMI_2")
    add_slide_layout_from_layout_name(prs, "BAPA_KAMI_3")




def add_preacher_page(prs, PASTOR_TITLE_ID, PASTOR_TITLE_DE, PASTOR_NAME):
    layout_name = "PREACHER" # renamed in the master template pptx file
    slide_layout = add_slide_layout_from_layout_name(prs, layout_name)

    print("preacher layout")
    # check_placeholders_in_slide(prs,slide_layout)

    de_preacher_placeholder = slide_layout.placeholders[10]
    id_preacher_placeholder = slide_layout.placeholders[11]

    de_preacher_placeholder.text = PASTOR_TITLE_DE + " " + PASTOR_NAME
    id_preacher_placeholder.text = PASTOR_TITLE_ID + " " + PASTOR_NAME

    
    

def add_appostle_creed_page(prs):
    #loop and search for layout with the name "0_APOSTLE_CREED_1", "1_APOSTLE_CREED_1" until "5_APOSTLE_CREED_1"
    for i in range(0,6):
        layout_name = str(i) + "_APOSTLE_CREED_1"
        slide_layout = add_slide_layout_from_layout_name(prs, layout_name)
        print( str(i) + " apostle creed layout")

def decide_offering_purpose_layout_name(next_sunday_date):
    # if next sunday is the first sunday of the month, then its "P_PENGINJILAN"
    # second sunday, then its "P_SEKOLAH"
    # third sunday, then its "P_MANDAT"
    # fourth sunday, then its "P_PEMBANGUNAN"
    # fifth sunday, then its "P_DIAKONIA"
    output = ""
    if next_sunday_date.day <= 7:
        output = "P_PENGINJILAN"
    elif next_sunday_date.day <= 14:
        output = "P_SEKOLAH"
    elif next_sunday_date.day <= 21:
        output = "P_MANDAT"
    elif next_sunday_date.day <= 28:
        output = "P_PEMBANGUNAN"
    else:
        output = "P_DIAKONIA"
    return output
    

def add_secondary_offering_purpose_page(prs, offering_purpose):
    # if offering_purpose is P_PENGINJILAN, then add slide with layout name "P_PENGINJILAN"
    # if offering_purpose is P_SEKOLAH, then add slide with layout name "P_SEKOLAH"
    layout_name = offering_purpose
    slide_layout = add_slide_layout_from_layout_name(prs, layout_name)




def add_doxology_page(prs):
    add_slide_layout_from_layout_name(prs, "0_DOXOLOGY")
    add_slide_layout_from_layout_name(prs, "1_DOXOLOGY")
    add_slide_layout_from_layout_name(prs, "2_DOXOLOGY")

def add_amen_page(prs):
    add_slide_layout_from_layout_name(prs, "3_DOXOLOGY")



def add_bekantmachung_page(prs):
    try: 
        for i in range(0,20): # 0,1,2,3,4,5,6,7
            layout_name = str(i) + "_WARTA"  

            slide_layout = add_slide_layout_from_layout_name(prs, layout_name)

            if i == 3:
                # preparing texts and placeholders
                texts=["Makan Malam & \n Persekutuan Doa", 
                       "Abendessen & Gebetkreis", 
                       "Setiap Jumat di Minggu Ganjil \n 18:30 ", 
                       "Freitags der ungeraden Woche \n 18:30 "]
                placeholders = get_placeholders_in_slide(prs, slide_layout)
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 4:
                # preparing texts and placeholders
                texts=["Pemahaman Alkitab", 
                       "Sabtu, 15:00  \n Berner Heerweg 60", 
                       "Bibelstunde", 
                       "Samstags, 15:00  \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 5:
                # preparing texts and placeholders
                texts=["Katekisasi Online", 
                       "Sabtu, setiap 2 minggu, 13:00 "]
                placeholders = get_placeholders_in_slide(prs, slide_layout)
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 6:
                # preparing texts and placeholders
                texts=["Master Class", 
                       "Setiap Sabtu minggu genap \n 14:00 "]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 7:
                # preparing texts and placeholders
                texts=["Katekisasi Nikah", 
                       "Pendaftaran: Nina"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 8:
                # preparing texts and placeholders
                texts=["Latihan Koor", 
                       "Minggu, 14:00 - 15:00 \n Berner Heerweg 60", 
                       "Chorübung", 
                       "Sonntags, 14:00 - 15:00 \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 9:
                # preparing texts and placeholders
                texts=["Persekutuan Doa", 
                       "Setiap Minggu, 15:30  \n Berner Heerweg 60", 
                       "Gebetkreis", 
                       "Sonntags, 15:30 \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1
            
            elif i == 10:
                # preparing texts and placeholders
                texts=["Ibadah Minggu", 
                       "Minggu, 14:00 \n Berner Heerweg 60", 
                       "Sonntagsgottesdienst", 
                       "Sonntag, 14:00  \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 11:
                # preparing texts and placeholders
                texts=["Ibadah Minggu", 
                       "Setiap Minggu, 16:00 \n Berner Heerweg 60", 
                       "Sonntagsgottesdienst", 
                       "Jeden Sonntag, 16:00  \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 12:
                # preparing texts and placeholders
                texts=["Ibadah Minggu & Perjamuan Kudus", 
                       "Minggu, 9:00 \n Berner Heerweg 60", 
                       "Sonntagsgottesdienst & Abendmahl", 
                       "Sonntag, 9:00 \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 13:
                # preparing texts and placeholders
                texts=["Persekutuan Doa", 
                       "Minggu, 11:00 \n Berner Heerweg 60", 
                       "Gebetkreis", 
                       "Sonntag, 11:00 \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 14:
                # preparing texts and placeholders
                texts=["Persekutuan Doa", 
                       "Setiap Minggu, 15:30 \n Berner Heerweg 60", 
                       "Gebetkreis", 
                       "Sonntags, 15:30 \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 15:
                # preparing texts and placeholders
                texts=["Sekolah Minggu", 
                       "Minggu, 9:00 \n Berner Heerweg 60", 
                       "Sonntagsschule", 
                       "Sonntag, 9:00 \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1

            elif i == 16:
                # preparing texts and placeholders
                texts=["Sekolah Minggu", 
                       "Setiap Minggu, 16:00 \n Berner Heerweg 60", 
                       "Sonntagsschule", 
                       "Jeden Sonntag, 16:00 \n Berner Heerweg 60"]
                placeholders = get_placeholders_in_slide(prs, slide_layout) 
                index = 0

                # filling the placeholders with the texts
                # the first text should be in the first placeholder, the second text should be in the second placeholder, and so on
                for placeholder in placeholders:
                    slide_layout.placeholders[placeholder].text = texts[index]
                    index += 1
            


    except IndexError:
        print("Invalid placeholder index. It's okay, we can continue.")


#################### IN MEMORY SONG SAVING #########################################


def insert_slides_from_google_drive_folder(prs, song_number):
    """
    prs: Presentation object
    song_slides: list of song_slide objects, ex: [song_slide1, song_slide2, song_slide3]
    """
    song_images_byte = []
    song_images = []
    
    song_images_byte = download_new_song_pipeline(song_number).values()
    song_images_byte = list(song_images_byte)
    make_song_slides_from_song_imgs(prs, song_images_byte)



def make_song_slides_from_song_imgs(prs, song_images_byte):
    """
    song_imgs: list of song_img objects, ex: [song_img1, song_img2, song_img3] (from Pujian.py)
    the somg_img object we get from google drives
    """
    song_slides = []
    for song_image_byte in song_images_byte:
        # create a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # add the picture to the slide
        pic = slide.shapes.add_picture(song_image_byte, Inches(0), Inches(0), height=prs.slide_height, width=prs.slide_width) # i dont need song_image_byte.getvalue() because we are using the BytesIO directly


def make_song_slide_from_song_img(prs, song_img):
    """
    song_img: song_img object, ex: song_img1 (from Pujian.py)
    the somg_img object we get from google drives
    """
    # create a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # add the picture to the slide
    pic = slide.shapes.add_picture(song_img, Inches(0), Inches(0), height=prs.slide_height, width=prs.slide_width)
    
    





######## Helper functions ####################################################

def check_placeholders_in_slide_index(prs, layout_index):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    print('in layout %d :' % layout_index)
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            print('id: %d, name: %s' % (phf.idx, phf.type))
    print('check done. Remember, the id is more important than the position')

def check_placeholders_in_slide(prs, slide):
    print('in slide %s' % slide)
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            print('id: %d, name: %s' % (phf.idx, phf.type))
    print('check done. Remember, the id is more important than the position')

def get_placeholders_in_slide(prs, slide):
    print('in slide %s' % slide)
    placeholders = []
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            placeholders.append(phf.idx)
    return placeholders

    

#### TEST FUNCTIONS ############################################################

def test_insert_slides_from_pict_folder(prs, folder_path):
    folder_path = os.path.join(os.path.dirname(__file__), 'Sample', '2', '2')
    insert_slides_from_pict_folder(prs, folder_path)
